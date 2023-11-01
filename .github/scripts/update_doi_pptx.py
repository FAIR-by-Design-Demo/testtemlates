from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import re
import glob
import yaml
import requests

DOI_PATTERN = r'\b10\.\d{4,}(?:\.\d+)*\/[^\s]+\b'
SANDBOX=False
ZENODO_URL = 'https://zenodo.org/api'

if ZENODO_SANDBOX.lower() == 'true':
    ZENODO_URL = 'https://sandbox.zenodo.org/api'

def replace_text_last_slide(presentation_path: str, new_doi: str):
    presentation = Presentation(presentation_path)
    last_slide = presentation.slides[-1]

    for shape in last_slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                doi_match = re.findall(DOI_PATTERN, paragraph.text)
                if 'To cite this presentation' in paragraph.text or (len(doi_match)>0 and doi_match[0]!=new_doi):
                    print(f'Found replacement candidate in {presentation_path}, replacing...')
                    paragraph.text = ''
                    run_instructions = paragraph.add_run()
                    run_instructions.text = 'Cite instructions on:\n'
                    run_instructions.font.name = 'Quicksand'
                    run_instructions.font.size = Pt(12)
                    run_instructions.font.color.rgb = RGBColor(0, 84, 144)
                    run_doi = paragraph.add_run()
                    run_doi.text = f'https://doi.org/{new_doi}'
                    run_doi.font.name = 'Quicksand'
                    run_doi.font.size = Pt(12)
                    run_doi.font.color.rgb = RGBColor(229, 115, 0)
                    presentation.save(presentation_path)
                    return

def parse_doi(doi: str):
    pattern = r'zenodo\.(\d+)'
    match = re.search(pattern, doi)
    if match:
        identifier = match.group(1)
        return identifier
    raise Exception('Invalid DOI.')

def get_zenodo_record(identifier: str):
    response = requests.get(f'{ZENODO_URL}/records/{identifier}')
    collection = response.json().get('conceptrecid')
    if response.status_code == 200 and collection is not None:
        return collection
    else:
        return ''

def get_universal_doi():
    with open('CITATION.cff', 'r') as yaml_file:
        yaml_content = yaml.safe_load(yaml_file)
        latest_doi = yaml_content.get('doi')
        return get_zenodo_record(latest_doi)

universal_doi = get_universal_doi()
print(f'Universal DOI: {universal_doi}')

for file in glob.glob('**/*.pptx', recursive=True):
    if 'venv/' in file:
        continue
    print(f'Checking {file}...')
    replace_text_last_slide(file, universal_doi)
